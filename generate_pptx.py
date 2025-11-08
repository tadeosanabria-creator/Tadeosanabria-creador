#!/usr/bin/env python3
"""
Genera infografia.pptx a partir de infografia_editable.svg (colores y diseño conservados).
Requisitos: cairosvg, python-pptx
Instalar: pip install cairosvg python-pptx
Ejecutar: python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches
import cairosvg
import os
import sys

SVG_FILE = "infografia_editable.svg"
PNG_FILE = "infografia.png"
PPTX_FILE = "infografia.pptx"

if not os.path.exists(SVG_FILE):
    sys.exit(f"ERROR: No se encontró '{SVG_FILE}' en la raíz del repositorio. Asegúrate de que exista.")

print(f"Convirtiendo {SVG_FILE} -> {PNG_FILE} (rasterizando SVG)...")
# Ajusta el tamaño del PNG si quieres más resolución (dpi_scale)
# Por defecto cairosvg generará una buena resolución; si quieres más, usa output_width/output_height
cairosvg.svg2png(url=SVG_FILE, write_to=PNG_FILE)

print("Creando PPTX y agregando la imagen en una sola diapositiva 16:9...")
prs = Presentation()
# Ajusta tamaño de la presentación a 16:9 (13.333 x 7.5 pulgadas)
prs.slide_width = Inches(13.3333333333)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]  # layout en blanco
slide = prs.slides.add_slide(blank_slide_layout)

# Inserta la imagen y la escala para llenar la diapositiva
left = top = Inches(0)
pic = slide.shapes.add_picture(PNG_FILE, left, top, width=prs.slide_width, height=prs.slide_height)

prs.save(PPTX_FILE)
print(f"Generado: {PPTX_FILE}")
print("Revisa el PPTX en PowerPoint para convertir la imagen en formas si deseas editar los elementos por separado.")
